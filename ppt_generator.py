#!/usr/bin/env python3
"""
PowerPoint Report Generator for Industry Analysis Framework
Automated PPT generation using python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import json
from datetime import datetime
from pathlib import Path
import logging

class IndustryAnalysisPPTGenerator:
    """Generates professional PowerPoint presentations from analysis results"""
    
    def __init__(self):
        self.prs = Presentation()
        self.setup_theme()
    
    def setup_theme(self):
        """Setup presentation theme and styling"""
        # Set slide size to widescreen
        self.prs.slide_width = Inches(16)
        self.prs.slide_height = Inches(9)
        
        # Define color scheme
        self.colors = {
            'primary': RGBColor(0, 51, 102),      # Dark blue
            'secondary': RGBColor(0, 153, 204),   # Light blue
            'accent': RGBColor(255, 102, 0),      # Orange
            'text': RGBColor(51, 51, 51),         # Dark gray
            'background': RGBColor(248, 249, 250)  # Light gray
        }
    
    def add_title_slide(self, analysis_date: str):
        """Add title slide"""
        slide_layout = self.prs.slide_layouts[0]  # Title slide layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = "Hong Kong Industry Analysis Framework"
        title.text_frame.paragraphs[0].font.size = Pt(44)
        title.text_frame.paragraphs[0].font.color.rgb = self.colors['primary']
        
        subtitle.text = f"Medical R&D and Patent Brokerage Industries\nAnalysis Date: {analysis_date}"
        subtitle.text_frame.paragraphs[0].font.size = Pt(24)
        subtitle.text_frame.paragraphs[0].font.color.rgb = self.colors['text']
    
    def add_executive_summary(self, results: dict):
        """Add executive summary slide"""
        slide_layout = self.prs.slide_layouts[1]  # Title and content layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "Executive Summary"
        title.text_frame.paragraphs[0].font.color.rgb = self.colors['primary']
        
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.clear()
        
        # Key findings
        summary = results['summary']
        
        p = tf.paragraphs[0]
        p.text = f"• Total companies analyzed: {summary['total_companies_analyzed']}"
        p.font.size = Pt(18)
        
        p = tf.add_paragraph()
        p.text = f"• Medical R&D companies identified: {summary['medical_rd_companies']}"
        p.font.size = Pt(18)
        
        p = tf.add_paragraph()
        p.text = f"• Patent brokerage companies identified: {summary['patent_brokerage_companies']}"
        p.font.size = Pt(18)
        
        p = tf.add_paragraph()
        p.text = "• Critical Gap: HSIC lacks dedicated codes for emerging industries"
        p.font.size = Pt(18)
        p.font.color.rgb = self.colors['accent']
        
        p = tf.add_paragraph()
        p.text = "• HK R&D expenditure: 0.99% of GDP vs Singapore's 1.89%"
        p.font.size = Pt(18)
    
    def add_classification_framework(self, results: dict):
        """Add industry classification framework slide"""
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "Industry Classification Framework (PPT Framework)"
        title.text_frame.paragraphs[0].font.color.rgb = self.colors['primary']
        
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.clear()
        
        classification = results['classification_framework']
        
        # Medical R&D section
        p = tf.paragraphs[0]
        p.text = "I. Medical R&D Industry"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = self.colors['secondary']
        
        medical_rd = classification['medical_rd']
        p = tf.add_paragraph()
        p.text = f"   a. ISIC: {medical_rd['isic_code']} ({medical_rd['description']})"
        p.font.size = Pt(16)
        p.level = 1
        
        p = tf.add_paragraph()
        p.text = f"   b. HSIC gap: {medical_rd['hsic_gap']}"
        p.font.size = Pt(16)
        p.level = 1
        
        p = tf.add_paragraph()
        p.text = "   c. Proposed framework:"
        p.font.size = Pt(16)
        p.level = 1
        
        for item in medical_rd['proposed_framework']:
            p = tf.add_paragraph()
            p.text = f"      - {item}"
            p.font.size = Pt(14)
            p.level = 2
        
        # Patent Brokerage section
        p = tf.add_paragraph()
        p.text = "\nII. Patent Brokerage"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = self.colors['secondary']
        
        patent_brokerage = classification['patent_brokerage']
        p = tf.add_paragraph()
        p.text = f"   a. ISIC: {patent_brokerage['isic_code']} ({patent_brokerage['description']})"
        p.font.size = Pt(16)
        p.level = 1
        
        p = tf.add_paragraph()
        p.text = f"   b. HSIC gap: {patent_brokerage['hsic_gap']}"
        p.font.size = Pt(16)
        p.level = 1
        
        p = tf.add_paragraph()
        p.text = "   c. Proposed framework:"
        p.font.size = Pt(16)
        p.level = 1
        
        for item in patent_brokerage['proposed_framework']:
            p = tf.add_paragraph()
            p.text = f"      - {item}"
            p.font.size = Pt(14)
            p.level = 2
    
    def add_company_listings(self, results: dict):
        """Add company listings slides"""
        # Medical R&D Companies
        if results['company_listings']['medical_rd']:
            slide_layout = self.prs.slide_layouts[1]
            slide = self.prs.slides.add_slide(slide_layout)
            
            title = slide.shapes.title
            title.text = "Medical R&D Companies Identified"
            title.text_frame.paragraphs[0].font.color.rgb = self.colors['primary']
            
            content = slide.placeholders[1]
            tf = content.text_frame
            tf.clear()
            
            for i, company in enumerate(results['company_listings']['medical_rd'], 1):
                p = tf.paragraphs[0] if i == 1 else tf.add_paragraph()
                p.text = f"{i}. {company['name']}"
                p.font.size = Pt(18)
                p.font.bold = True
                p.font.color.rgb = self.colors['secondary']
                
                p = tf.add_paragraph()
                p.text = f"   Nature: {company['business_nature']}"
                p.font.size = Pt(14)
                p.level = 1
                
                p = tf.add_paragraph()
                p.text = f"   Location: {company['address']}"
                p.font.size = Pt(14)
                p.level = 1
                
                p = tf.add_paragraph()
                p.text = f"   ISIC/HSIC: {company['isic_code']}/{company['hsic_code']}"
                p.font.size = Pt(14)
                p.level = 1
        
        # Patent Brokerage Companies
        if results['company_listings']['patent_brokerage']:
            slide_layout = self.prs.slide_layouts[1]
            slide = self.prs.slides.add_slide(slide_layout)
            
            title = slide.shapes.title
            title.text = "Patent Brokerage Companies Identified"
            title.text_frame.paragraphs[0].font.color.rgb = self.colors['primary']
            
            content = slide.placeholders[1]
            tf = content.text_frame
            tf.clear()
            
            for i, company in enumerate(results['company_listings']['patent_brokerage'], 1):
                p = tf.paragraphs[0] if i == 1 else tf.add_paragraph()
                p.text = f"{i}. {company['name']}"
                p.font.size = Pt(18)
                p.font.bold = True
                p.font.color.rgb = self.colors['secondary']
                
                p = tf.add_paragraph()
                p.text = f"   Nature: {company['business_nature']}"
                p.font.size = Pt(14)
                p.level = 1
                
                p = tf.add_paragraph()
                p.text = f"   Location: {company['address']}"
                p.font.size = Pt(14)
                p.level = 1
                
                p = tf.add_paragraph()
                p.text = f"   ISIC/HSIC: {company['isic_code']}/{company['hsic_code']}"
                p.font.size = Pt(14)
                p.level = 1
    
    def add_market_gaps_analysis(self, results: dict):
        """Add market gaps analysis slide"""
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "Industry Gap Analysis"
        title.text_frame.paragraphs[0].font.color.rgb = self.colors['primary']
        
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.clear()
        
        market_analysis = results['market_analysis']
        
        # Barriers section
        p = tf.paragraphs[0]
        p.text = "Key Development Barriers:"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = self.colors['accent']
        
        # Medical R&D barriers
        p = tf.add_paragraph()
        p.text = "Medical R&D:"
        p.font.size = Pt(16)
        p.font.bold = True
        p.level = 1
        
        for barrier in market_analysis['barriers']['medical_rd']:
            p = tf.add_paragraph()
            p.text = f"• {barrier}"
            p.font.size = Pt(14)
            p.level = 2
        
        # Patent brokerage barriers
        p = tf.add_paragraph()
        p.text = "Patent Brokerage:"
        p.font.size = Pt(16)
        p.font.bold = True
        p.level = 1
        
        for barrier in market_analysis['barriers']['patent_brokerage']:
            p = tf.add_paragraph()
            p.text = f"• {barrier}"
            p.font.size = Pt(14)
            p.level = 2
    
    def add_opportunities_slide(self, results: dict):
        """Add opportunities and recommendations slide"""
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "Development Opportunities & Recommendations"
        title.text_frame.paragraphs[0].font.color.rgb = self.colors['primary']
        
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.clear()
        
        market_analysis = results['market_analysis']
        
        # Opportunities section
        p = tf.paragraphs[0]
        p.text = "Strategic Opportunities:"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = self.colors['secondary']
        
        # Medical R&D opportunities
        p = tf.add_paragraph()
        p.text = "Medical R&D:"
        p.font.size = Pt(16)
        p.font.bold = True
        p.level = 1
        
        for opportunity in market_analysis['opportunities']['medical_rd']:
            p = tf.add_paragraph()
            p.text = f"• {opportunity}"
            p.font.size = Pt(14)
            p.level = 2
        
        # Patent brokerage opportunities
        p = tf.add_paragraph()
        p.text = "Patent Brokerage:"
        p.font.size = Pt(16)
        p.font.bold = True
        p.level = 1
        
        for opportunity in market_analysis['opportunities']['patent_brokerage']:
            p = tf.add_paragraph()
            p.text = f"• {opportunity}"
            p.font.size = Pt(14)
            p.level = 2
    
    def add_key_findings(self, results: dict):
        """Add key findings slide"""
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "Key Findings & Critical Gaps"
        title.text_frame.paragraphs[0].font.color.rgb = self.colors['primary']
        
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.clear()
        
        key_findings = results['key_findings']
        
        p = tf.paragraphs[0]
        p.text = "Industry Findings:"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = self.colors['secondary']
        
        p = tf.add_paragraph()
        p.text = f"• Medical R&D: {key_findings['medical_rd']}"
        p.font.size = Pt(16)
        
        p = tf.add_paragraph()
        p.text = f"• Patent Brokerage: {key_findings['patent_brokerage']}"
        p.font.size = Pt(16)
        
        p = tf.add_paragraph()
        p.text = "\nCritical Classification Gaps:"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = self.colors['accent']
        
        for gap in key_findings['critical_gaps']:
            p = tf.add_paragraph()
            p.text = f"• {gap}"
            p.font.size = Pt(16)
        
        p = tf.add_paragraph()
        p.text = f"\nDevelopment Potential:"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = self.colors['secondary']
        
        p = tf.add_paragraph()
        p.text = key_findings['development_potential']
        p.font.size = Pt(16)
    
    def add_chart_slide(self, chart_path: str, title: str):
        """Add slide with chart image"""
        slide_layout = self.prs.slide_layouts[6]  # Blank layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Add title
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(15), Inches(1))
        title_frame = title_shape.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(32)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = self.colors['primary']
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Add chart image if exists
        if Path(chart_path).exists():
            slide.shapes.add_picture(chart_path, Inches(1), Inches(2), Inches(14), Inches(6))
    
    def generate_presentation(self, results: dict, output_file: str = None) -> str:
        """Generate complete PowerPoint presentation"""
        if not output_file:
            timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
            output_file = f'HK_Industry_Analysis_{timestamp}.pptx'
        
        logging.info("Generating PowerPoint presentation...")
        
        # Add all slides
        analysis_date = datetime.fromisoformat(results['summary']['analysis_date']).strftime('%B %d, %Y')
        
        self.add_title_slide(analysis_date)
        self.add_executive_summary(results)
        self.add_classification_framework(results)
        self.add_company_listings(results)
        self.add_market_gaps_analysis(results)
        self.add_opportunities_slide(results)
        self.add_key_findings(results)
        
        # Add chart slides if charts exist
        chart_dir = Path('charts')
        if chart_dir.exists():
            rd_chart = chart_dir / 'rd_expenditure_comparison.png'
            patent_chart = chart_dir / 'patent_trends.png'
            
            if rd_chart.exists():
                self.add_chart_slide(str(rd_chart), 'R&D Expenditure Comparison: Hong Kong vs Singapore')
            
            if patent_chart.exists():
                self.add_chart_slide(str(patent_chart), 'Medical R&D Patent Filings & IP Services Growth')
        
        # Save presentation
        self.prs.save(output_file)
        logging.info(f"PowerPoint presentation saved as: {output_file}")
        
        return output_file

def generate_ppt_from_json(json_file: str, output_file: str = None) -> str:
    """Generate PowerPoint presentation from JSON results file"""
    with open(json_file, 'r', encoding='utf-8') as f:
        results = json.load(f)
    
    generator = IndustryAnalysisPPTGenerator()
    return generator.generate_presentation(results, output_file)

if __name__ == "__main__":
    # Example usage
    import sys
    
    if len(sys.argv) > 1:
        json_file = sys.argv[1]
        output_file = sys.argv[2] if len(sys.argv) > 2 else None
        ppt_file = generate_ppt_from_json(json_file, output_file)
        print(f"PowerPoint presentation generated: {ppt_file}")
    else:
        print("Usage: python ppt_generator.py <json_results_file> [output_file.pptx]")